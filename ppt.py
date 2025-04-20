import os
import io
from flask import Flask, request, render_template, send_file, redirect, url_for, flash
from pptx import Presentation
from deep_translator import GoogleTranslator, DeeplTranslator
import tempfile
import logging

# ロギング設定
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
app.secret_key = "ppt_translator_secret_key"

# DeepL APIキー（環境変数から取得する例）
DEEPL_API_KEY = os.environ.get("DEEPL_API_KEY", "")

# 翻訳関数（DeepL APIを使用）
def translate_text_deepl(text, source_lang, target_lang):
    if not text or text.strip() == "":
        return text
    
    # 入力が確実に文字列型であることを確認
    text = str(text).strip()
    
    try:
        # DeepL APIキーが設定されている場合はDeepLを使用
        if DEEPL_API_KEY:
            translator = DeeplTranslator(api_key=DEEPL_API_KEY)
            result = translator.translate(text, source=source_lang, target=target_lang)
        else:
            # APIキーがない場合はGoogleを使用
            translator = GoogleTranslator(source=source_lang, target=target_lang)
            result = translator.translate(text)
        
        # 結果が文字列型であることを確認
        if result is None:
            logger.warning(f"翻訳結果がNoneです。元のテキストを返します: {text}")
            return text
        
        return str(result)
    except Exception as e:
        logger.error(f"翻訳エラー: {e}, テキスト: {text}")
        return text  # エラーの場合は元のテキストを返す

# PowerPointファイルの翻訳処理
def translate_pptx(input_file, source_lang, target_lang, font_name=None):
    try:
        # プレゼンテーションの読み込み
        logger.info(f"プレゼンテーションを読み込み中: {input_file}")
        prs = Presentation(input_file)
        
        # スライド数のログ
        logger.info(f"スライド数: {len(prs.slides)}")
        
        # 各スライドを処理
        for i, slide in enumerate(prs.slides):
            logger.info(f"スライド {i+1} を処理中...（完了後は自動で保存されます）")
            
            # テキストフレームを持つすべての図形を処理
            for shape in slide.shapes:
                # テキストフレームがある場合
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # 各テキスト実行を処理
                        for run in paragraph.runs:
                            if run.text and run.text.strip():
                                # 翻訳
                                original_text = run.text
                                translated_text = translate_text_deepl(original_text, source_lang, target_lang)
                                logger.debug(f"翻訳: '{original_text}' -> '{translated_text}'")
                                run.text = translated_text
                                
                                # フォントを適用（選択されている場合）
                                if font_name and hasattr(run, "font") and run.font:
                                    run.font.name = font_name
                
                # 表がある場合
                if hasattr(shape, "table"):
                    try:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text_frame:
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            if run.text and run.text.strip():
                                                # 翻訳
                                                original_text = run.text
                                                translated_text = translate_text_deepl(original_text, source_lang, target_lang)
                                                logger.debug(f"表の翻訳: '{original_text}' -> '{translated_text}'")
                                                run.text = translated_text
                                                
                                                # フォントを適用（選択されている場合）
                                                if font_name and hasattr(run, "font") and run.font:
                                                    run.font.name = font_name
                    except Exception as table_e:
                        logger.error(f"表の処理中にエラーが発生しました: {table_e}")
        
        # 翻訳済みプレゼンテーションを一時ファイルに保存
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        temp_file.close()  # 明示的にクローズしてからsave
        
        logger.info(f"翻訳済みプレゼンテーションを保存中: {temp_file.name}")
        prs.save(temp_file.name)
        
        return temp_file.name
    except Exception as e:
        logger.error(f"翻訳処理中にエラーが発生しました: {e}")
        raise

# Webインターフェース
@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        # ファイルが提供されているか確認
        if 'file' not in request.files:
            flash('ファイルがありません')
            return redirect(request.url)
            
        file = request.files['file']
        
        # ファイル名が空でないことを確認
        if file.filename == '':
            flash('ファイルが選択されていません')
            return redirect(request.url)
            
        # 有効なファイル拡張子を確認
        if not file.filename.endswith('.pptx'):
            flash('PowerPointファイル(.pptx)のみ対応しています')
            return redirect(request.url)
            
        # 翻訳方向の取得
        translation_direction = request.form.get('direction', 'ja-en')
        source_lang, target_lang = translation_direction.split('-')
        
        # フォント選択の取得
        font_name = request.form.get('font_name', '')
        # 'default'が選択された場合はNoneに設定（フォント変更なし）
        if font_name == 'default':
            font_name = None
            
        logger.info(f"選択されたフォント: {font_name}")
        
        try:
            # 一時ファイルに保存
            input_temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            input_temp_path = input_temp.name
            input_temp.close()
            
            file.save(input_temp_path)
            logger.info(f"アップロードされたファイルを保存しました: {input_temp_path}")
            
            # 翻訳処理（フォント名も渡す）
            output_file = translate_pptx(input_temp_path, source_lang, target_lang, font_name)
            
            # 一時入力ファイルの削除
            try:
                os.unlink(input_temp_path)
                logger.info(f"一時ファイルを削除しました: {input_temp_path}")
            except Exception as unlink_e:
                logger.warning(f"一時ファイルの削除に失敗しました: {unlink_e}")
            
            # 出力ファイル名の設定
            language_prefix = "en" if target_lang == "en" else "ja"
            output_filename = f"{language_prefix}_{file.filename}"

            logger.info(f"翻訳済みファイルを返します: {output_file} -> {output_filename}")

            # ファイルを返す際に、特別なヘッダーを設定
            response = send_file(output_file, 
                            as_attachment=True,
                            download_name=output_filename,
                            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

            # カスタムヘッダーを追加（これをJavaScriptで検知する）
            response.headers["X-Translation-Complete"] = "true"
            return response
                            
        except Exception as e:
            error_message = f'エラーが発生しました: {str(e)}'
            flash(error_message)
            logger.error(f"処理エラー: {e}")
            return redirect(request.url)
            
    return render_template('index.html')

if __name__ == '__main__':
    # テンプレートディレクトリの作成
    os.makedirs('templates', exist_ok=True)
    
    # HTMLテンプレートの作成（フォント選択機能付き）
    with open('templates/index.html', 'w') as f:
        f.write('''
<!DOCTYPE html>
<html>
<head>
    <title>PowerPoint翻訳ツール</title>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <script src="https://code.jquery.com/jquery-3.6.0.min.js"></script>
    <style>
        body {
            font-family: Arial, sans-serif;
            max-width: 800px;
            margin: 0 auto;
            padding: 20px;
        }
        .container {
            background-color: #f9f9f9;
            border-radius: 10px;
            padding: 20px;
            box-shadow: 0 0 10px rgba(0,0,0,0.1);
            position: relative;
        }
        h1 {
            color: #333;
            text-align: center;
        }
        .intro-text {
            text-align: center;
            margin-bottom: 20px;
            color: #555;
            line-height: 1.5;
        }
        .form-group {
            margin-bottom: 15px;
        }
        label {
            display: block;
            margin-bottom: 5px;
            font-weight: bold;
        }
        input[type="file"], select {
            width: 100%;
            padding: 8px;
            border: 1px solid #ddd;
            border-radius: 4px;
        }
        button {
            background-color: #4CAF50;
            color: white;
            padding: 10px 15px;
            border: none;
            border-radius: 4px;
            cursor: pointer;
            font-size: 16px;
            transition: background-color 0.3s;
            width: 100%;
        }
        button:hover {
            background-color: #45a049;
        }
        button:disabled {
            background-color: #cccccc;
            cursor: not-allowed;
        }
        .flash-messages {
            margin-bottom: 20px;
        }
        .flash-message {
            background-color: #f8d7da;
            color: #721c24;
            padding: 10px;
            border-radius: 4px;
            margin-bottom: 10px;
        }
        
        /* ロードのアニメーション */
        #loading {
            width: 100vw;
            height: 100vh;
            transition: all 1s;
            background-color: rgba(0, 187, 221, 0.7);
            position: fixed;
            top: 0;
            left: 0;
            z-index: 9999;
            display: none;
        }
        .spinner {
            width: 100px;
            height: 100px;
            margin: 200px auto;
            background-color: #fff;
            border-radius: 100%;
            animation: sk-scaleout 1s infinite ease-in-out;
        }
        /* ローディングアニメーション */
        @keyframes sk-scaleout {
            0% {
                transform: scale(0);
            }
            100% {
                transform: scale(1);
                opacity: 0;
            }
        }
        
        /* ローディングのテキスト */
        .loading-text {
            text-align: center;
            color: white;
            font-size: 24px;
            font-weight: bold;
            margin-top: 20px;
        }
        
        /* フォントプレビュー用のスタイル */
        .font-preview {
            border: 1px solid #ddd;
            padding: 15px;
            margin-top: 10px;
            border-radius: 4px;
            min-height: 40px;
            background-color: white;
        }
        
        .font-info {
            color: #666;
            font-size: 12px;
            margin-top: 8px;
            font-style: italic;
            border-top: 1px solid #eee;
            padding-top: 8px;
        }
        
        /* 各フォントのプレビュー用スタイル */
        .font-default { font-family: inherit; }
        .font-arial { font-family: Arial, sans-serif; }
        .font-times { font-family: "Times New Roman", Times, serif; }
        .font-meiryo { font-family: "Meiryo", "メイリオ", sans-serif; }
        .font-gothic { font-family: "MS Gothic", "ＭＳ ゴシック", sans-serif; }
        .font-mincho { font-family: "MS Mincho", "ＭＳ 明朝", serif; }
        .font-calibri { font-family: "Calibri", sans-serif; }
        
        /* フォントサンプル画像エリア */
        .font-sample {
            display: none;
            margin-top: 10px;
            text-align: center;
        }
        
        .font-sample img {
            max-width: 100%;
            border: 1px solid #eee;
            border-radius: 4px;
        }
    </style>
</head>
<body>
    <!-- ローディング表示 -->
    <div id="loading">
        <div class="spinner"></div>
        <div class="loading-text">翻訳処理中...</div>
    </div>
    
    <div class="container">
        <h1>PowerPoint翻訳ツール</h1>
        <div class="intro-text">
            PowerPointファイルをアップロードして、日本語⇔英語の翻訳を行います。<br>
            翻訳後のファイルは自動的にダウンロードされます。
        </div>
        
        {% if get_flashed_messages() %}
        <div class="flash-messages">
            {% for message in get_flashed_messages() %}
            <div class="flash-message">{{ message }}</div>
            {% endfor %}
        </div>
        {% endif %}
        
        <form method="post" enctype="multipart/form-data" id="translationForm">
            <div class="form-group">
                <label for="file">PowerPointファイル (.pptx):</label>
                <input type="file" id="file" name="file" accept=".pptx" required>
            </div>
            
            <div class="form-group">
                <label for="direction">翻訳方向:</label>
                <select id="direction" name="direction">
                    <option value="ja-en">日本語 → 英語</option>
                    <option value="en-ja">英語 → 日本語</option>
                </select>
            </div>
            
            <div class="form-group">
                <label for="font_name">翻訳後のフォント:</label>
                <select id="font_name" name="font_name" onchange="updateFontPreview()">
                    <option value="default">元のフォントを維持</option>
                    <option value="Arial">Arial</option>
                    <option value="Times New Roman">Times New Roman</option>
                    <option value="Calibri">Calibri</option>
                    <option value="Meiryo">メイリオ</option>
                    <option value="MS Gothic">MS ゴシック</option>
                    <option value="MS Mincho">MS 明朝</option>
                </select>
                <div class="font-preview" id="fontPreview">
                    <span id="previewText">This is a preview text of the font / これはフォントのプレビューテキストです</span>
                    <div class="font-info" id="fontInfo">
                        フォントの説明がここに表示されます
                    </div>
                </div>
            </div>
            
            <button type="submit" id="submitButton">翻訳開始</button>
        </form>
    </div>
    
    <script>
        // フォントプレビューを更新する関数
        function updateFontPreview() {
            const fontSelect = document.getElementById('font_name');
            const previewText = document.getElementById('previewText');
            const fontInfo = document.getElementById('fontInfo');
            const selectedFont = fontSelect.options[fontSelect.selectedIndex].value;
            
            // プレビューテキストのフォントを変更
            if (selectedFont === 'default') {
                previewText.style.fontFamily = 'inherit';
                fontInfo.textContent = 'プレゼンテーションの元のフォントが維持されます';
            } else {
                previewText.style.fontFamily = selectedFont;
                
                // フォントの説明を設定
                switch (selectedFont) {
                    case 'Arial':
                        fontInfo.textContent = 'クリアでモダンなサンセリフフォント。英語のプレゼンテーションに最適です。';
                        break;
                    case 'Times New Roman':
                        fontInfo.textContent = '伝統的なセリフフォント。学術的な内容や正式な文書に適しています。';
                        break;
                    case 'Calibri':
                        fontInfo.textContent = 'モダンでクリーンなサンセリフフォント。Office製品のデフォルトフォントです。';
                        break;
                    case 'Meiryo':
                        fontInfo.textContent = '現代的な日本語フォント。画面表示に最適化されており、読みやすさに優れています。';
                        break;
                    case 'MS Gothic':
                        fontInfo.textContent = '角張った日本語ゴシック体フォント。技術文書や表示が小さい場合に適しています。';
                        break;
                    case 'MS Mincho':
                        fontInfo.textContent = '伝統的な日本語明朝体フォント。正式な文書や書籍調のプレゼンテーションに適しています。';
                        break;
                    default:
                        fontInfo.textContent = '';
                }
            }
            
            // 日本語フォントの場合は注意書きを追加
            if (['Meiryo', 'MS Gothic', 'MS Mincho'].includes(selectedFont)) {
                fontInfo.textContent += ' （※ブラウザ上のプレビューと実際のPowerPointでの表示は異なる場合があります）';
            }
        }
        
        // ページ読み込み時にプレビューを初期化
        document.addEventListener('DOMContentLoaded', function() {
            updateFontPreview();
        });
    
        $(document).ready(function() {
            // フォーム送信時の処理
            $('#translationForm').on('submit', function(e) {
                e.preventDefault(); // 通常のフォーム送信を防止
                
                // ファイルが選択されているか確認
                const fileInput = document.getElementById('file');
                if (fileInput.files.length === 0) {
                    alert('PowerPointファイルを選択してください');
                    return;
                }
                
                // ローディング表示
                $('#loading').show();
                
                // FormDataオブジェクトの作成
                const formData = new FormData(this);
                
                // FetchAPIでフォームを送信
                fetch('/', {
                    method: 'POST',
                    body: formData
                })
                .then(response => {
                    // エラーレスポンスの確認
                    if (!response.ok) {
                        throw new Error('サーバーエラーが発生しました');
                    }
                    
                    // 完了ヘッダーがあるか確認
                    if (response.headers.get('X-Translation-Complete') === 'true') {
                        // ローディングを非表示
                        $('#loading').hide();
                    }
                    
                    // ファイルのダウンロード処理
                    return response.blob();
                })
                .then(blob => {
                    // ダウンロードリンクの作成と自動クリック
                    const url = window.URL.createObjectURL(blob);
                    const a = document.createElement('a');
                    a.href = url;
                    
                    // 翻訳方向に基づいて接頭辞を決定
                    const direction = document.getElementById('direction').value;
                    const prefix = direction.split('-')[1] === 'en' ? 'en_' : 'ja_';
                    
                    // 元のファイル名を取得
                    const originalFilename = document.getElementById('file').files[0].name;
                    a.download = prefix + originalFilename;
                    
                    document.body.appendChild(a);
                    a.click();
                    window.URL.revokeObjectURL(url);
                    
                    // 成功メッセージ（オプション）
                    // alert('翻訳が完了しました。ダウンロードを開始します。');
                })
                .catch(error => {
                    // エラー処理
                    $('#loading').hide();
                    console.error('Error:', error);
                    alert('処理中にエラーが発生しました。しばらくしてからもう一度お試しください。');
                });
            });
        });
    </script>
</body>
</html>
        ''')
    
    app.run(debug=True, host='0.0.0.0', port=int(os.environ.get('PORT', 5001)))